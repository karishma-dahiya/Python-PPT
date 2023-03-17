from wand.image import Image
import glob 
from pptx import Presentation 
from pptx.util import Inches 




directory = './img'


for file in glob.iglob(f'{directory}/*'):
    with Image(filename=file) as image:
        image.transform(resize='x400')
        with Image(filename='./logo.png') as water:
            water.transform(resize='x40')
            with image.clone() as watermark:
                watermark.watermark(water,0,1,1)
                watermark.save(filename=file)



prs = Presentation()

for file in glob.iglob(f'{directory}/*'):
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[1])


    title_shape = slide.shapes.title
    placeholder = slide.placeholders[1]

    title_shape.text = 'Title'

    placeholder.text = 'Subtitle'
    top = Inches(2.5) 
    left = Inches(1)
    image = slide.shapes.add_picture(file,left,top)

    
    
prs.save('test.pptx')                


